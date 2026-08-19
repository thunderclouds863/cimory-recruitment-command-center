from __future__ import annotations
from io import BytesIO
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from core.config import PPT_TEMPLATE, BRAND_LOGO, BRAND_BLUE, BRAND_RED, TEXT_DARK, TEXT_MUTED
from core.metrics import kpis, category_summary, recruiter_summary, directorate_summary, weekly_trend, aging_summary, funnel_summary

BLUE=RGBColor.from_string(BRAND_BLUE); RED=RGBColor.from_string(BRAND_RED); DARK=RGBColor.from_string(TEXT_DARK); MUTED=RGBColor.from_string(TEXT_MUTED); WHITE=RGBColor(255,255,255); LIGHT=RGBColor(246,247,251)

def _clean_template(prs):
    ids=list(prs.slides._sldIdLst)
    for sid in ids:
        rId=sid.rId; prs.part.drop_rel(rId); prs.slides._sldIdLst.remove(sid)

def _slide(prs,title,subtitle=None):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=WHITE
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(.12)).fill.solid(); slide.shapes[-1].fill.fore_color.rgb=BLUE; slide.shapes[-1].line.fill.background()
    if BRAND_LOGO.exists(): slide.shapes.add_picture(str(BRAND_LOGO),Inches(.35),Inches(.2),height=Inches(.55))
    tb=slide.shapes.add_textbox(Inches(.45),Inches(.92),Inches(12.1),Inches(.55)); p=tb.text_frame.paragraphs[0]; p.text=title; p.font.name="Raleway"; p.font.size=Pt(24); p.font.bold=True; p.font.color.rgb=DARK
    if subtitle:
        sb=slide.shapes.add_textbox(Inches(.47),Inches(1.46),Inches(12),Inches(.32)); p=sb.text_frame.paragraphs[0]; p.text=subtitle; p.font.name="Raleway"; p.font.size=Pt(10); p.font.color.rgb=MUTED
    return slide

def _metric(slide,x,y,w,h,value,label,color=BLUE):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=LIGHT; sh.line.color.rgb=RGBColor(226,232,240)
    tf=sh.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=str(value); p.alignment=PP_ALIGN.CENTER; p.font.name="Raleway"; p.font.size=Pt(25); p.font.bold=True; p.font.color.rgb=color
    p=tf.add_paragraph(); p.text=label.upper(); p.alignment=PP_ALIGN.CENTER; p.font.name="Raleway"; p.font.size=Pt(8); p.font.color.rgb=MUTED

def _table(slide,df,x,y,w,h,font=8):
    if df is None or df.empty:return
    df=df.copy().fillna(""); rows=min(len(df)+1,14); df=df.head(rows-1)
    t=slide.shapes.add_table(len(df)+1,len(df.columns),Inches(x),Inches(y),Inches(w),Inches(h)).table
    for j,c in enumerate(df.columns):
        cell=t.cell(0,j); cell.text=str(c); cell.fill.solid(); cell.fill.fore_color.rgb=BLUE
        for p in cell.text_frame.paragraphs: p.font.color.rgb=WHITE; p.font.bold=True; p.font.name="Raleway"; p.font.size=Pt(font)
    for i,(_,r) in enumerate(df.iterrows(),1):
        for j,c in enumerate(df.columns):
            val=r[c]; val=f"{val:.1f}%" if isinstance(val,float) and ("Rate" in c or "SLA" in c) else val
            t.cell(i,j).text=str(val)
            for p in t.cell(i,j).text_frame.paragraphs: p.font.name="Raleway"; p.font.size=Pt(font-1); p.font.color.rgb=DARK

def _bar_chart(slide,df,cat_col,val_col,x,y,w,h,title):
    if df.empty:return
    d=df.head(12).copy(); data=ChartData(); data.categories=[str(v) for v in d[cat_col]]; data.add_series(val_col,[float(v) for v in d[val_col]])
    chart=slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,Inches(x),Inches(y),Inches(w),Inches(h),data).chart; chart.has_title=True; chart.chart_title.text_frame.text=title; chart.has_legend=False
    chart.value_axis.has_major_gridlines=True

def _line_chart(slide,df,x,y,w,h):
    if df.empty:return
    d=df.tail(20); data=ChartData(); data.categories=[str(v) for v in d["Week"]]; data.add_series("FPTK Processed",[int(v) for v in d["FPTK Processed"]]); data.add_series("Closed",[int(v) for v in d["Closed"]])
    chart=slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,Inches(x),Inches(y),Inches(w),Inches(h),data).chart; chart.has_title=True; chart.chart_title.text_frame.text="Weekly FPTK Inflow vs Closure"; chart.has_legend=True; chart.legend.position=XL_LEGEND_POSITION.BOTTOM

def build_ceo_update_pptx(fptk_df,user=None,title_date=None):
    prs=Presentation(str(PPT_TEMPLATE)) if PPT_TEMPLATE.exists() else Presentation(); _clean_template(prs); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    m=kpis(fptk_df); cat=category_summary(fptk_df); rec=recruiter_summary(fptk_df); direc=directorate_summary(fptk_df); wk=weekly_trend(fptk_df); age=aging_summary(fptk_df); funnel=funnel_summary(user)
    s=_slide(prs,"TALENT ACQUISITION UPDATE",title_date or "Auto-generated from Recruitment Command Center")
    if BRAND_LOGO.exists(): s.shapes.add_picture(str(BRAND_LOGO),Inches(4.1),Inches(2.15),width=Inches(5.2))
    s=_slide(prs,"Executive Summary — Seluruh FPTK"); labels=[("FPTK MASUK",m['total']),("OPEN",m['open']),("CLOSED",m['closed']),("CANCEL",m['cancel']),("FILL RATE",f"{m['fill_rate']:.2f}%"),("CLOSE SESUAI SLA",f"{m['closed_sla']:.2f}%")]
    for i,(lab,val) in enumerate(labels): _metric(s,.55+i*2.05,2.0,1.78,1.3,val,lab,BLUE if lab not in {"OPEN","CANCEL"} else RED)
    _table(s,cat[[c for c in ["Category","FPTK","Open","Closed","Cancel","Fill Rate","Closed sesuai SLA"] if c in cat]],.6,3.65,12.1,2.65,8)
    s=_slide(prs,"FPTK by Directorate","Where demand is concentrated"); _bar_chart(s,direc,"Directorate","Total",.7,1.8,12,4.9,"Total FPTK by Directorate")
    s=_slide(prs,"Recruiter Workload & Performance"); _bar_chart(s,rec,"Recruiter","Open",.7,1.8,6.1,4.9,"Open FPTK by Recruiter"); _table(s,rec[["Recruiter","Total","Open","Closed","Over SLA","Fill Rate"]],7.05,1.8,5.7,4.9,8)
    s=_slide(prs,"Recruitment Funnel"); _bar_chart(s,funnel.sort_values("Count",ascending=False) if not funnel.empty else funnel,"Stage","Count",.7,1.8,12,4.9,"Candidates by Passed Stage")
    s=_slide(prs,"Open FPTK Aging & SLA Risk"); _bar_chart(s,age,"Aging","Open",.7,1.8,5.7,4.9,"Open FPTK Aging"); risk=fptk_df[(fptk_df.status=="OP")].copy(); risk["Days to SLA"]=(risk.deadline_sla-pd.Timestamp.today().normalize()).dt.days; risk=risk.sort_values("Days to SLA").head(12); _table(s,risk[["position","pic_recruiter","level_fptk","deadline_sla","Days to SLA"]],6.7,1.8,6.0,4.9,7)
    s=_slide(prs,"Pemenuhan SDM — Weekly Trend"); _line_chart(s,wk,.7,1.8,12,4.9)
    for cat_name in ["Manager","Level 3","Level 2C Below","STO","CLAP & FGDP"]:
        d=fptk_df[fptk_df.filter_category.fillna("")==cat_name]
        if d.empty: continue
        mm=kpis(d); s=_slide(prs,f"Executive Summary — {cat_name}")
        for i,(lab,val) in enumerate([("FPTK",mm['total']),("OPEN",mm['open']),("CLOSED",mm['closed']),("CANCEL",mm['cancel']),("FILL RATE",f"{mm['fill_rate']:.2f}%"),("CLOSE SESUAI SLA",f"{mm['closed_sla']:.2f}%")]): _metric(s,.55+i*2.05,1.8,1.78,1.25,val,lab,BLUE if lab not in {"OPEN","CANCEL"} else RED)
        open_d=d[d.status=="OP"].copy().sort_values("fptk_date"); _table(s,open_d[["fptk_date","position","business_unit","level_fptk","pic_recruiter","deadline_sla"]],.6,3.45,12.1,2.9,7)
    out=BytesIO(); prs.save(out); return out.getvalue()
